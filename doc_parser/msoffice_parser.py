"""Parser for Microsoft Office documents."""

# Standard library imports
import tempfile  # For creating temporary files and directories
from typing import List  # For type hints

# Third-party library imports
from langchain_community.document_loaders import UnstructuredExcelLoader  # For loading unstructured Excel documents
from langchain_community.document_loaders import Docx2txtLoader  # For loading text from Word documents
from pptx import Presentation  # For working with PowerPoint presentations
from io import BytesIO  # For working with in-memory binary data

# Local imports
from .parserbase import DocumentParser

class MsDocParser(DocumentParser):
    """Parser for Microsoft Word documents."""

    def __init__(self, document) -> None:
        """
        Initialize the MsDocParser.

        Parameters:
            document: bytes
                The binary data of the Word document.
        """
        super().__init__(document)

    def extract_raw_text(self) -> List[str]:
        """
        Extract raw text from the Word document.

        Returns:
            List[str]:
                A list of strings representing the extracted text from the document.
        """
        with tempfile.NamedTemporaryFile(mode='wb', delete=False) as tmpf:
            # Write buffer to temporary file.
            tmpf.write(self.buffer)
            tmpf.seek(0)
            # Load document.
            loader = Docx2txtLoader(tmpf.name)
            data = loader.load()
            # Convert list of 'Document' to list of string.
            return [d.page_content for d in data]


class MsExcelParser(DocumentParser):
    """Parser for Microsoft Excel documents."""

    def __init__(self, document) -> None:
        """
        Initialize the MsExcelParser.

        Parameters:
            document: bytes
                The binary data of the Excel document.
        """
        super().__init__(document)

    def extract_raw_text(self) -> List[str]:
        """
        Extract raw text from the Excel document.

        Returns:
            List[str]:
                A list of strings representing the extracted text from the document.
        """
        with tempfile.NamedTemporaryFile(mode='wb', delete=False) as tmpf:
            # Write buffer to temporary file.
            tmpf.write(self.buffer)
            tmpf.seek(0)
            # Load document.
            loader = UnstructuredExcelLoader(tmpf.name, mode='elements')
            data = loader.load()
            # Convert list of 'Document' to list of string.
            return [d.page_content for d in data]  

class MsPptParser(DocumentParser):
    """Parser for Microsoft PowerPoint documents."""

    def __init__(self, document) -> None:
        """
        Initialize the MsPptParser.

        Parameters:
            document: bytes
                The binary data of the PowerPoint document.
        """
        super().__init__(document)

    def extract_raw_text(self) -> List[str]:
        """
        Extract raw text from the PowerPoint document.

        Returns:
            List[str]:
                A list of strings representing the extracted text from the document.
        """
        parser = Presentation(BytesIO(self.buffer))
        extracted_text = []

        # Extract text from slides.
        for slides in parser.slides:
            for shape in slides.shapes:
                if hasattr(shape, "text"):
                    extracted_text.append(shape.text)

        return extracted_text
